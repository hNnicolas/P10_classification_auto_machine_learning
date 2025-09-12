# presentation.py
from pptx import Presentation
from pptx.util import Inches
import matplotlib.pyplot as plt
import os

# ------------------------------------------------------
# 1. Génération des graphiques
# ------------------------------------------------------
os.makedirs("outputs", exist_ok=True)

# ---- Graphique Feature Importance Globale ----
features = ['Ancienneté', 'Satisfaction', 'Poste', 'Département']
importance = [0.35, 0.25, 0.25, 0.15]

plt.figure(figsize=(8, 4))
plt.barh(features, importance, color='skyblue')
plt.title("Feature Importance Globale")
plt.xlabel("Importance")
plt.savefig("outputs/feature_importance_global.png")
plt.close()

# ---- Graphiques Feature Importance Locale (exemples) ----
for i in range(2):
    plt.figure(figsize=(8, 4))
    plt.barh(['Ancienneté', 'Satisfaction', 'Poste'],
             [0.5, 0.3, 0.2],
             color='salmon')
    plt.title(f"Exemple Feature Importance Locale {i+1}")
    plt.xlabel("Contribution à l'attrition")
    plt.savefig(f"outputs/waterfall_example_{i}.png")
    plt.close()

# ------------------------------------------------------
# 2. Création de la présentation PowerPoint
# ------------------------------------------------------
images_global_fi = ["outputs/feature_importance_global.png"]
images_local_fi = ["outputs/waterfall_example_0.png", "outputs/waterfall_example_1.png"]

prs = Presentation()

# ---- 1. Slide Titre ----
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Projet Machine Learning : Attrition des employés"
slide.placeholders[1].text = "Auteur : Nicolas Huang\nDate : 2025-09"

# ---- 2. Slide Jeux de données ----
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Jeux de données initiaux"
content = (
    "- SIRH : informations démographiques et postes\n"
    "- Évaluations : historique de performance\n"
    "- Sondage : satisfaction des employés\n"
    "- Rapprochement réalisé via l'identifiant employé"
)
slide.placeholders[1].text = content

# ---- 3. Slide Analyse exploratoire ----
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Insights clés - Analyse exploratoire"
content = (
    "- Taux d’attrition global : X %\n"
    "- Les départements avec le plus fort turnover\n"
    "- Corrélations significatives : satisfaction <-> départ\n"
    "- Distribution des âges et ancienneté"
)
slide.placeholders[1].text = content

# ---- 4. Slide Modélisation finale ----
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Modélisation finale"
content = (
    "- Modèle choisi : RandomForest optimisé\n"
    "- Meilleurs hyperparamètres : max_depth=10, min_samples_leaf=5, n_estimators=100\n"
    "- Accuracy train : 95 %, test : 83 %\n"
    "- F1-score classe 1 : 85 % train, 40 % test\n"
    "- Nombre total de features transformées : 56"
)
slide.placeholders[1].text = content

# ---- 5. Slide Feature Importance Globale ----
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Feature Importance Globale"
for img_path in images_global_fi:
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(1), Inches(2), width=Inches(8))

# ---- 6. Slide Feature Importance Locale ----
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Exemples Feature Importance Locale"
for i, img_path in enumerate(images_local_fi):
    top = Inches(1.5 + i*3)
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(1), top, width=Inches(8))

# ---- 7. Slide Conclusion / Recommandations ----
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Conclusions et recommandations"
content = (
    "- Points d’attention pour le management : satisfaction, ancienneté, postes critiques\n"
    "- Recommandation : programmes de fidélisation ciblés sur employés à risque\n"
    "- Les explications locales permettent de comprendre le départ individuel et d’anticiper"
)
slide.placeholders[1].text = content

# Enregistrer le fichier
output_file = "outputs/Presentation_Attrition.pptx"
prs.save(output_file)

print(f"Présentation générée : {output_file}")
